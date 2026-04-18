from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class Brand:
    bg: RGBColor = RGBColor(0x08, 0x08, 0x08)  # #080808
    card: RGBColor = RGBColor(0x11, 0x11, 0x11)  # #111111
    accent: RGBColor = RGBColor(0x39, 0xD3, 0x53)  # #39D353
    text: RGBColor = RGBColor(0xF0, 0xED, 0xE8)  # #F0EDE8
    highlight: RGBColor = RGBColor(0xE8, 0xF5, 0x5A)  # #E8F55A


BRAND = Brand()


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text_box(slide, x, y, w, h, text: str, *, size=16, bold=False, color=BRAND.text, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _card(slide, x, y, w, h, *, fill=BRAND.card, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = (line if line else fill)
    shape.line.width = Pt(1)
    return shape


def _logo_and_footer(slide, slide_no: int, logo_path: Path):
    # Logo
    slide.shapes.add_picture(str(logo_path), Inches(0.45), Inches(0.28), height=Inches(0.42))

    # Slide number
    _text_box(
        slide,
        Inches(12.9),
        Inches(7.05),
        Inches(0.6),
        Inches(0.3),
        f"{slide_no:02d}",
        size=12,
        bold=True,
        color=RGBColor(0xAA, 0xAA, 0xAA),
        align=PP_ALIGN.RIGHT,
    )


def _takeaway(slide, text: str):
    _text_box(
        slide,
        Inches(0.7),
        Inches(7.0),
        Inches(12.0),
        Inches(0.45),
        text,
        size=14,
        bold=True,
        color=BRAND.accent,
        align=PP_ALIGN.LEFT,
    )


def _headline(slide, text: str):
    return _text_box(slide, Inches(0.7), Inches(0.9), Inches(12.6), Inches(0.9), text, size=44, bold=True)


def _source(slide, text: str):
    _text_box(
        slide,
        Inches(0.7),
        Inches(6.65),
        Inches(10.5),
        Inches(0.35),
        text,
        size=10,
        bold=False,
        color=RGBColor(0x99, 0x99, 0x99),
    )


def _big_number(slide, x, y, number: str, label: str, *, color=BRAND.text):
    _text_box(slide, x, y, Inches(4.0), Inches(0.8), number, size=64, bold=True, color=color)
    _text_box(slide, x, y + Inches(0.8), Inches(4.0), Inches(0.5), label, size=16, bold=False, color=RGBColor(0xC8, 0xC8, 0xC8))


def _arrow(slide, x, y, w, h, text: str):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BRAND.card
    shp.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = BRAND.text
    return shp


def _flow_column(slide, x, y, w, title: str, steps: list[str]):
    _text_box(slide, x, y, w, Inches(0.35), title, size=16, bold=True, color=BRAND.highlight)
    cur_y = y + Inches(0.5)
    box_h = Inches(0.5)
    gap = Inches(0.18)
    for i, step in enumerate(steps):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, cur_y, w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = BRAND.card
        shp.line.color.rgb = RGBColor(0x2B, 0x2B, 0x2B)
        shp.line.width = Pt(1)
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{i+1}. {step}"
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.color.rgb = BRAND.text
        if i < len(steps) - 1:
            slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, x + w / 2 - Inches(0.12), cur_y + box_h, Inches(0.24), gap)
        cur_y += box_h + gap


def build_deck(out_path: Path, *, mark_logo: Path, wordmark_logo: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    def add_slide(n: int):
        s = prs.slides.add_slide(blank)
        _set_bg(s, BRAND.bg)
        if n != 1:
            _logo_and_footer(s, n, mark_logo)
        return s

    # Slide 1 — Cover
    s1 = add_slide(1)
    _set_bg(s1, BRAND.bg)
    s1.shapes.add_picture(str(wordmark_logo), Inches(4.1), Inches(2.0), width=Inches(5.2))
    _text_box(s1, Inches(2.4), Inches(3.2), Inches(8.6), Inches(0.4), "Smart Proximity Automatic Recognizing Operator", size=18, bold=False, align=PP_ALIGN.CENTER)
    _text_box(s1, Inches(2.1), Inches(3.75), Inches(9.2), Inches(0.6), "\"Just arrive. SPAROO handles the rest.\"", size=22, bold=True, color=BRAND.accent, align=PP_ALIGN.CENTER)
    _text_box(s1, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.5), "Nithin Reddy Jeeru | Founder    •    ANITS, Visakhapatnam | 2025", size=14, bold=False, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)

    # Slide 2 — Problem
    s2 = add_slide(2)
    _headline(s2, "India’s UPI payment process\nis broken at the last meter.")
    ex = _card(s2, Inches(0.7), Inches(2.35), Inches(6.5), Inches(2.05), fill=BRAND.card)
    _text_box(
        s2,
        Inches(1.05),
        Inches(2.55),
        Inches(5.9),
        Inches(1.6),
        "It is 9 PM. Auto ride ends.\nDriver says — UPI kar do bhai.\nQR code is faded.\nUPI ID mistyped twice.\n2 minutes wasted.\nEvery. Single. Ride.",
        size=16,
        bold=False,
        color=BRAND.text,
    )
    _big_number(s2, Inches(8.0), Inches(2.4), "640M+", "UPI transactions / day")
    _big_number(s2, Inches(8.0), Inches(3.85), "6Cr+", "MSMEs (mostly micro)")
    _big_number(s2, Inches(8.0), Inches(5.3), "≈3 min", "avg payment friction (field observation)")
    _takeaway(s2, "The payment technology works. The discovery problem is unsolved.")
    _source(s2, "Sources: PIB (UPI >640M/day), Economic Survey 2025–26 (UPI stats), Ministry of MSME/Udyam (enterprise count).")

    # Slide 3 — Solution
    s3 = add_slide(3)
    _headline(s3, "SPAROO eliminates the discovery problem.")
    _text_box(
        s3,
        Inches(0.7),
        Inches(2.15),
        Inches(7.0),
        Inches(3.4),
        "A small device sits inside the auto.\nWhen you come near — your phone\nautomatically shows the driver’s details.\nName. Photo. UPI ID. Verified.\n\nYou enter the amount. Tap pay.\nDone. Under 10 seconds.\n\nNo QR scanning.\nNo typing UPI IDs.\nNo asking.\n\nTelugu voice confirmation:\nవంద రూపాయలు క్రెడిట్ అయింది.",
        size=18,
    )
    _text_box(s3, Inches(8.1), Inches(2.15), Inches(4.8), Inches(0.5), "Simple 4-step", size=16, bold=True, color=BRAND.highlight)
    _arrow(s3, Inches(8.1), Inches(2.8), Inches(4.8), Inches(0.65), "Device broadcasts")
    _arrow(s3, Inches(8.1), Inches(3.6), Inches(4.8), Inches(0.65), "Phone detects")
    _arrow(s3, Inches(8.1), Inches(4.4), Inches(4.8), Inches(0.65), "Details appear")
    _arrow(s3, Inches(8.1), Inches(5.2), Inches(4.8), Inches(0.65), "Pay done")
    _takeaway(s3, "From 2–3 minutes to under 10 seconds — without changing how UPI works.")

    # Slide 4 — USP
    s4 = add_slide(4)
    _headline(s4, "SPAROO is a never-before solution\nfor India’s informal economy.")
    cols_x = [Inches(0.7), Inches(4.7), Inches(8.7)]
    titles = ["USP 1 — Auto Detection", "USP 2 — Telugu Voice", "USP 3 — Two-Layer Security"]
    bodies = [
        "No other payment solution\nautomatically detects the merchant.\nZero user action required.",
        "Like a soundbox — but at ₹98.\nTelugu voice plays every payment.\nDriver hears it instantly.",
        "Device stores zero sensitive data.\nUPI ID lives only in secure cloud.\nKYC verified merchants (Razorpay).",
    ]
    for x, t, b in zip(cols_x, titles, bodies):
        _card(s4, x, Inches(2.5), Inches(3.7), Inches(3.0))
        _text_box(s4, x + Inches(0.35), Inches(2.75), Inches(3.05), Inches(0.55), t, size=16, bold=True, color=BRAND.accent)
        _text_box(s4, x + Inches(0.35), Inches(3.35), Inches(3.05), Inches(2.0), b, size=16)
    _takeaway(s4, "Built for auto drivers, street vendors, kirana shops — the real India.")

    # Slide 5 — Competition
    s5 = add_slide(5)
    _headline(s5, "Every alternative fails India’s informal economy.")
    _card(s5, Inches(0.7), Inches(2.3), Inches(8.1), Inches(3.6))
    rows = [
        ("Auto-detect", "❌", "❌", "✅"),
        ("Works for autos", "⚠️", "❌", "✅"),
        ("Device cost", "₹20", "₹500+", "₹98"),
        ("Offline capable", "✅", "❌", "✅"),
        ("Voice confirm", "❌", "❌", "✅"),
        ("Merchant verified", "❌", "⚠️", "✅"),
    ]
    headers = ["Feature", "QR Code", "ToneTag", "SPAROO"]
    colx = [Inches(1.0), Inches(4.1), Inches(5.7), Inches(7.25)]
    _text_box(s5, colx[0], Inches(2.45), Inches(3.0), Inches(0.35), headers[0], size=14, bold=True, color=BRAND.highlight)
    _text_box(s5, colx[1], Inches(2.45), Inches(1.4), Inches(0.35), headers[1], size=14, bold=True, color=BRAND.text, align=PP_ALIGN.CENTER)
    _text_box(s5, colx[2], Inches(2.45), Inches(1.4), Inches(0.35), headers[2], size=14, bold=True, color=BRAND.text, align=PP_ALIGN.CENTER)
    _text_box(s5, colx[3], Inches(2.45), Inches(1.6), Inches(0.35), headers[3], size=14, bold=True, color=BRAND.accent, align=PP_ALIGN.CENTER)
    y = Inches(2.85)
    for feat, qr, tone, spa in rows:
        _text_box(s5, colx[0], y, Inches(3.0), Inches(0.35), feat, size=14, color=RGBColor(0xD0, 0xD0, 0xD0))
        _text_box(s5, colx[1], y, Inches(1.4), Inches(0.35), qr, size=14, bold=True, align=PP_ALIGN.CENTER)
        _text_box(s5, colx[2], y, Inches(1.4), Inches(0.35), tone, size=14, bold=True, align=PP_ALIGN.CENTER)
        _text_box(s5, colx[3], y, Inches(1.6), Inches(0.35), spa, size=14, bold=True, color=BRAND.accent, align=PP_ALIGN.CENTER)
        y += Inches(0.45)
    _card(s5, Inches(9.1), Inches(2.3), Inches(3.5), Inches(3.6))
    _text_box(s5, Inches(9.45), Inches(2.55), Inches(2.9), Inches(0.6), "Barrier to entry", size=16, bold=True, color=BRAND.highlight)
    _text_box(
        s5,
        Inches(9.45),
        Inches(3.1),
        Inches(2.9),
        Inches(2.7),
        "ToneTag raised ~$78M.\nThat validates the market.\n\nBut they target enterprise.\nWe own the informal segment.\n\nHardware + software +\nmerchant network + KYC\n= ~18 months to copy.\n\nBy then: 10,000 drivers.",
        size=14,
    )
    _takeaway(s5, "SPAROO wins where QR and enterprise NFC can’t: fast discovery in the real world.")

    # Slide 6 — Revenue
    s6 = add_slide(6)
    _headline(s6, "Three streams. All growing together.")
    _card(s6, Inches(0.7), Inches(2.35), Inches(12.0), Inches(2.6))
    _text_box(s6, Inches(1.1), Inches(2.6), Inches(3.7), Inches(2.2), "1) Device Sale\n₹200–300 price\n₹98 cost at scale\n~3× margin", size=18, bold=True, color=BRAND.text)
    _text_box(s6, Inches(5.0), Inches(2.6), Inches(3.7), Inches(2.2), "2) Transaction Fee\n0.3% per payment\n10L tx/month → ₹3L/month\nScales automatically", size=18, bold=True, color=BRAND.text)
    _text_box(s6, Inches(8.9), Inches(2.6), Inches(3.6), Inches(2.2), "3) B2B SaaS\n₹50/driver/month\n10k drivers → ₹5L/month\n~86% margin", size=18, bold=True, color=BRAND.text)
    _card(s6, Inches(0.7), Inches(5.25), Inches(6.2), Inches(1.45))
    _text_box(s6, Inches(1.05), Inches(5.45), Inches(5.6), Inches(1.0), "3-year projection\nYear 1: ₹30L (pilot)\nYear 2: ₹1.5Cr (3 cities)\nYear 3: ₹8Cr (national)", size=18, bold=True, color=BRAND.accent)
    _card(s6, Inches(7.15), Inches(5.25), Inches(5.55), Inches(1.45))
    _text_box(s6, Inches(7.5), Inches(5.45), Inches(4.9), Inches(1.0), "Top cash burn\n1) Device manufacturing\n2) App development\n3) Merchant acquisition\n4) Operations", size=16, bold=False)
    _takeaway(s6, "Hardware starts revenue. Software makes it compounding.")

    # Slide 7 — Market
    s7 = add_slide(7)
    _headline(s7, "India’s micro-businesses.\nCompletely underserved.")
    _big_number(s7, Inches(0.9), Inches(2.35), "6.3Cr+", "MSMEs (majority micro)")
    _big_number(s7, Inches(0.9), Inches(3.85), "3Cr+", "Auto drivers (target beachhead)")
    _big_number(s7, Inches(0.9), Inches(5.35), "1.5Cr+", "Street vendors (next)")
    _card(s7, Inches(6.2), Inches(2.35), Inches(6.5), Inches(3.9))
    _text_box(s7, Inches(6.55), Inches(2.65), Inches(5.8), Inches(0.45), "Go-to-market channels", size=18, bold=True, color=BRAND.highlight)
    _text_box(
        s7,
        Inches(6.55),
        Inches(3.2),
        Inches(5.8),
        Inches(3.0),
        "1) Direct driver outreach (pilot)\n2) Auto unions + associations\n3) Fleet company B2B deals\n4) Driver word-of-mouth",
        size=18,
        bold=False,
    )
    _takeaway(s7, "Start with autos. Expand to every micro-merchant who can’t rely on a QR.")
    _source(s7, "Sources: Ministry of MSME Annual Report / Udyam portal registrations; segmentation = internal target prioritization.")

    # Slide 8 — Product + App flows
    s8 = add_slide(8)
    _headline(s8, "Two apps. One device.\nA complete payment ecosystem.")
    _card(s8, Inches(0.7), Inches(2.25), Inches(3.9), Inches(3.95))
    _text_box(s8, Inches(1.05), Inches(2.5), Inches(3.25), Inches(0.45), "SPAROO Beacon", size=18, bold=True, color=BRAND.accent)
    _text_box(
        s8,
        Inches(1.05),
        Inches(3.0),
        Inches(3.35),
        Inches(3.0),
        "Small BLE device\nSits in auto\nBroadcasts identity\nBattery: ~12 months\nCost at scale: ₹98\nTelugu voice confirm",
        size=16,
    )
    _card(s8, Inches(4.9), Inches(2.25), Inches(7.8), Inches(3.95))
    _flow_column(
        s8,
        Inches(5.25),
        Inches(2.5),
        Inches(3.6),
        "App 1: SPAROO (Customer)",
        [
            "Home shows Nearby merchants",
            "Start Scanning (BLE + Location)",
            "Tap Pay on merchant card",
            "Enter amount",
            "Razorpay checkout",
            "Success → saved to history",
        ],
    )
    _flow_column(
        s8,
        Inches(9.1),
        Inches(2.5),
        Inches(3.3),
        "App 2: SPAROO Partner",
        [
            "Login (Phone → OTP)",
            "Onboarding (Info + Bank)",
            "Link device ID",
            "Razorpay KYC (Route)",
            "Dashboard (earnings + status)",
            "Transactions + payouts",
        ],
    )
    _takeaway(s8, "One tiny beacon makes payments discoverable — apps make it scalable.")
    _source(s8, "Tech: Android (Kotlin) + BLE 5.0 + Supabase + Razorpay. (Implementation choice; not a market claim.)")

    # Slide 9 — Milestones
    s9 = add_slide(9)
    _headline(s9, "Built. Tested. Ready to scale.")
    _card(s9, Inches(0.7), Inches(2.3), Inches(5.9), Inches(3.9))
    _text_box(
        s9,
        Inches(1.05),
        Inches(2.55),
        Inches(5.3),
        Inches(3.4),
        "✅ BLE device prototype — working\n✅ Android customer app — complete\n✅ Partner app — built\n✅ Supabase backend — live\n✅ Real UPI payments — tested\n✅ Razorpay integration — complete\n✅ Telugu voice on device — working\n✅ Demo video — recorded\n\nStarted as RYDENT → failed → rebuilt as SPAROO",
        size=16,
    )
    _card(s9, Inches(6.9), Inches(2.3), Inches(5.8), Inches(1.85))
    _text_box(s9, Inches(7.25), Inches(2.55), Inches(5.2), Inches(1.35), "Next 12 months\nM1–2: 10-driver pilot (Vizag)\nM3–4: Register company + 30 devices\nM5–8: Raise funding + 50+ devices\nM9–12: Revenue live + Hyderabad", size=16, bold=True, color=BRAND.accent)
    _card(s9, Inches(6.9), Inches(4.35), Inches(5.8), Inches(1.85))
    _text_box(s9, Inches(7.25), Inches(4.6), Inches(5.2), Inches(1.35), "Next 3 years\nYear 2: 3 cities, 10k merchants\nYear 3: National, 1L merchants, B2B SaaS", size=16, bold=False)
    _takeaway(s9, "If funded today: 50 devices ordered in Week 1. First revenue in Month 2.")

    # Slide 10 — Funding
    s10 = add_slide(10)
    _headline(s10, "Zero external funding.\n100% student-built.")
    _big_number(s10, Inches(0.9), Inches(2.35), "₹0", "Funding raised")
    _big_number(s10, Inches(0.9), Inches(3.85), "5", "ANITS student team")
    _big_number(s10, Inches(0.9), Inches(5.35), "6+ months", "Time invested")
    _card(s10, Inches(6.3), Inches(2.35), Inches(6.4), Inches(3.9))
    _text_box(s10, Inches(6.65), Inches(2.65), Inches(5.7), Inches(0.5), "Asking: ₹5,00,000 seed", size=22, bold=True, color=BRAND.highlight)
    _text_box(
        s10,
        Inches(6.65),
        Inches(3.25),
        Inches(5.7),
        Inches(2.7),
        "Use of funds\n40% → 50-device pilot\n25% → custom PCB manufacturing\n20% → app + backend\n15% → marketing + ops\n\nOutcome\nPilot data in 60 days\nRevenue in 90 days\nNext round in 12 months",
        size=16,
    )
    _takeaway(s10, "We don’t need money to start. We need money to scale fast and measure right.")

    # Slide 11 — Team
    s11 = add_slide(11)
    _headline(s11, "Five students. One mission.\nZero excuses.")
    cards = [
        ("Nithin Reddy Jeeru", "Founder & Developer", "Built: Android app, BLE firmware,\nbackend, prototype hardware\n\n“I saw the problem daily.\nI refused to accept it.”"),
        ("Member 2", "Hardware & PCB", "Device design\nPCB roadmap\nManufacturing plan"),
        ("Member 3", "UI/UX Design", "App design\nBrand identity\nSPAROO visual language"),
        ("Member 4", "Business Development", "Driver outreach\nMarket research\nFleet partnerships"),
        ("Member 5", "Operations & Research", "Pilot planning\nData collection\nExpansion strategy"),
    ]
    x = Inches(0.7)
    y = Inches(2.25)
    w = Inches(4.0)
    h = Inches(1.7)
    for i, (name, role, body) in enumerate(cards):
        _card(s11, x, y, w, h)
        _text_box(s11, x + Inches(0.3), y + Inches(0.22), w - Inches(0.6), Inches(0.35), name, size=16, bold=True, color=BRAND.accent)
        _text_box(s11, x + Inches(0.3), y + Inches(0.55), w - Inches(0.6), Inches(0.28), role, size=12, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC))
        _text_box(s11, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), h - Inches(1.0), body, size=12)
        x += Inches(4.25)
        if i == 2:
            x = Inches(0.7)
            y += Inches(1.95)
    _card(s11, Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.55), fill=RGBColor(0x0D, 0x0D, 0x0D))
    _text_box(
        s11,
        Inches(1.0),
        Inches(6.48),
        Inches(11.4),
        Inches(0.35),
        "We started with RYDENT — it failed. We rebuilt as SPAROO — it works. A team that learns fast is the only team worth betting on.",
        size=14,
        bold=True,
        color=BRAND.highlight,
        align=PP_ALIGN.CENTER,
    )
    _takeaway(s11, "We are builders first — product, hardware, and field execution.")

    # Slide 12 — Closing
    s12 = add_slide(12)
    _headline(s12, "As long as there is that fire —\nyou will find a way.")
    _text_box(s12, Inches(0.7), Inches(2.25), Inches(12.6), Inches(0.4), "— Ratan N. Tata", size=18, bold=False, color=RGBColor(0xAA, 0xAA, 0xAA))
    _card(s12, Inches(0.7), Inches(3.05), Inches(12.0), Inches(2.5))
    _text_box(
        s12,
        Inches(1.1),
        Inches(3.35),
        Inches(11.2),
        Inches(2.0),
        "We found the problem.\nWe built the solution.\nWe failed once, learned, rebuilt.\n\n640M+ transactions daily.\n6+ crore micro-businesses waiting.\nZero smart discovery solutions for them.\n\nWe are not waiting for permission.\n\nJust arrive.\nSPAROO handles the rest.",
        size=20,
        bold=True,
        color=BRAND.text,
    )
    _text_box(s12, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.55), "Contact: Nithin Reddy Jeeru • Founder, SPAROO • ANITS, Visakhapatnam • sparoo.in", size=14, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    _source(s12, "Sources: PIB (UPI daily volume), Economic Survey 2025–26 Table 3.6 (UPI volume/value), DFS/NPCI (Mar 2026 peak month).")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    assets = here / "assets"
    out = here / "SPAROO_Pitch_Deck.pptx"
    build_deck(out, mark_logo=assets / "sparoo_mark.png", wordmark_logo=assets / "sparoo_wordmark.png")
    print(f"Wrote: {out}")

